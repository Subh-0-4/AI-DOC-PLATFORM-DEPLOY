from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
import os
from docx import Document
from pptx import Presentation

from ..auth import get_current_user, get_db
from .. import models

router = APIRouter(prefix="/export", tags=["export"])


# ---------- Helper to load project ----------
def _get_project_or_404(project_id: int, db: Session, current_user: models.User):
    project = (
        db.query(models.Project)
        .filter(models.Project.id == project_id, models.Project.owner_id == current_user.id)
        .first()
    )
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    return project


# ---------- DOCX Export ----------
@router.get("/docx/{project_id}")
def export_project_docx(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    project = _get_project_or_404(project_id, db, current_user)

    filename = f"project_{project.id}.docx"
    filepath = os.path.join("exports", filename)

    os.makedirs("exports", exist_ok=True)

    doc = Document()

    doc.add_heading(project.name, level=1)
    doc.add_paragraph(f"Main Topic: {project.main_topic}")
    doc.add_paragraph("\n")

    for section in project.sections:
        doc.add_heading(section.title, level=2)
        doc.add_paragraph(section.content or "")

    doc.save(filepath)
    return FileResponse(filepath, filename=filename)


# ---------- PPTX Export ----------
@router.get("/pptx/{project_id}")
def export_project_pptx(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    project = _get_project_or_404(project_id, db, current_user)

    filename = f"project_{project.id}.pptx"
    filepath = os.path.join("exports", filename)

    os.makedirs("exports", exist_ok=True)

    ppt = Presentation()

    # Title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = project.name
    subtitle.text = project.main_topic

    # Content slides
    for section in project.sections:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = section.title
        body.text = section.content or ""

    ppt.save(filepath)
    return FileResponse(filepath, filename=filename)
